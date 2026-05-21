from pathlib import Path


def parse_pptx(file_path: str | Path) -> str:
    """Extract text from a .pptx file using python-pptx."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    from pptx import Presentation
    prs = Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            slides_text.append("\n".join(slide_texts))
    return "\n\n".join(slides_text)
