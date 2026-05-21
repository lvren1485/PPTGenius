"""Main PPT generator: converts a structured outline into a professional .pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ..models.outline import PresentationOutline
from .templates import get_template
from .slides import (
    build_title_slide,
    build_section_slide,
    build_content_slide,
    build_two_column_slide,
    build_ending_slide,
)
from ..config import config


def generate_presentation(
    outline: PresentationOutline,
    template_name: str = "professional-blue",
    output_path: str | None = None,
) -> str:
    """Generate a professional .pptx file from a structured outline.

    Uses proper PPT slide layouts (selectable in PowerPoint's Layout picker).

    Args:
        outline: PresentationOutline with slides list.
        template_name: Name of the template to use.
        output_path: Output file path (auto-generated if None).

    Returns:
        Path to the generated .pptx file.
    """
    template = get_template(template_name)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = outline.slides
    total_slides = len(slides_data)

    for idx, slide_data in enumerate(slides_data):
        layout = slide_data.layout_type
        page_info = (idx, total_slides)

        if layout == "title":
            build_title_slide(prs, template, slide_data.title,
                              slide_data.subtitle or "", page_info)

        elif layout == "section":
            build_section_slide(prs, template, slide_data.title, page_info)

        elif layout == "two_column":
            mid = len(slide_data.bullets) // 2
            if mid < 1:
                mid = len(slide_data.bullets)
            build_two_column_slide(
                prs, template, slide_data.title,
                left_items=slide_data.bullets[:mid] or ["(left)"],
                right_items=slide_data.bullets[mid:] or ["(right)"],
                page_info=page_info,
            )

        elif layout == "ending":
            build_ending_slide(prs, template, slide_data.title,
                               slide_data.subtitle or "", page_info)

        else:
            build_content_slide(prs, template, slide_data.title,
                                slide_data.bullets, page_info)

    # Save
    if output_path is None:
        output_path = str(config.OUTPUT_DIR / "presentation.pptx")

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return str(path)
