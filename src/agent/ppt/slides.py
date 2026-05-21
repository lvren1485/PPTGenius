"""Slide builder functions — uses proper PPT layouts (selectable in PowerPoint)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .templates import Template
from .styles import (
    MARGIN_LEFT, CONTENT_WIDTH, TITLE_TOP, BODY_TOP,
)

# Layout indices in default python-pptx template
LAYOUT_TITLE = 0       # Title Slide
LAYOUT_CONTENT = 1     # Title and Content
LAYOUT_SECTION = 2     # Section Header
LAYOUT_TWO_COL = 3     # Two Content
LAYOUT_BLANK = 6       # Blank


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor. Strips alpha if 8 chars."""
    clean = hex_color.lstrip("#")
    if len(clean) >= 8:
        clean = clean[:6]
    r, g, b = bytes.fromhex(clean)
    return RGBColor(r, g, b)


def _set_bg(slide, color: str):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color)


def _add_shape(slide, left, top, width, height, fill_color=None):
    """Add a rectangle shape for visual decoration."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    else:
        shape.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                 font_size=18, color="#333333", bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = _hex_to_rgb(color)
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(4)
    return tf


def _add_bullets(tf, bullets, font_name, font_size, color):
    """Add bullet points to an existing text frame."""
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgb(color)
        p.level = 0
        p.space_after = Pt(8)
        p.space_before = Pt(4)
    return tf


def _add_page_number(slide, total, index, template: Template):
    """Add page number at bottom-right."""
    _add_textbox(
        slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
        f"{index + 1} / {total}",
        template.fonts.body, 10, template.colors.highlight,
        alignment=PP_ALIGN.RIGHT,
    )


def build_title_slide(prs: Presentation, template: Template, title: str,
                      subtitle: str = "", page_info: tuple = (1, 1)):
    """Build a professional title slide using Title Slide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _set_bg(slide, template.colors.bg)

    # Top accent bar
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08),
               template.colors.accent)

    # Title
    _add_textbox(
        slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
        title, template.fonts.title, 44, template.colors.accent,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
            subtitle, template.fonts.body, 22, template.colors.text,
            alignment=PP_ALIGN.CENTER,
        )

    # Bottom accent line
    _add_shape(slide, Inches(4), Inches(5.0), Inches(5.333), Inches(0.04),
               template.colors.highlight)

    # Date area
    from datetime import date
    _add_textbox(
        slide, Inches(1.5), Inches(6.5), Inches(5), Inches(0.4),
        date.today().strftime("%B %Y"),
        template.fonts.body, 12, template.colors.highlight,
        alignment=PP_ALIGN.LEFT,
    )

    _add_page_number(slide, page_info[1], page_info[0], template)
    return slide


def build_section_slide(prs: Presentation, template: Template, title: str,
                        page_info: tuple = (1, 1)):
    """Build a section divider slide using Section Header layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    _set_bg(slide, template.colors.accent)

    # Decorative left bar
    _add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
               template.colors.highlight)

    # Section number (automatic)
    _add_textbox(
        slide, Inches(1.5), Inches(1.5), Inches(3), Inches(0.6),
        f"SECTION {page_info[0]}",
        template.fonts.body, 14, template.colors.highlight,
        bold=True,
    )

    # Title
    _add_textbox(
        slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2),
        title, template.fonts.title, 40, template.colors.bg,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    _add_page_number(slide, page_info[1], page_info[0], template)
    return slide


def build_content_slide(prs: Presentation, template: Template, title: str,
                        bullets: list[str], page_info: tuple = (1, 1)):
    """Build a content slide using Title and Content layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    _set_bg(slide, template.colors.bg)

    # Top accent line
    _add_shape(slide, MARGIN_LEFT, Inches(0.08), Inches(2.5), Inches(0.05),
               template.colors.accent)

    # Title
    _add_textbox(
        slide, MARGIN_LEFT, TITLE_TOP, CONTENT_WIDTH, Inches(0.8),
        title, template.fonts.title, 30, template.colors.accent, bold=True,
    )

    # Accent underline
    _add_shape(slide, MARGIN_LEFT, Inches(1.2), Inches(1.5), Inches(0.03),
               template.colors.highlight)

    # Content
    if bullets:
        tf = _add_textbox(
            slide, MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, Inches(5.0),
            "", template.fonts.body, 18, template.colors.text,
        )
        tf.paragraphs[0].text = bullets[0]
        if len(bullets) > 1:
            _add_bullets(tf, bullets[1:], template.fonts.body, 18, template.colors.text)
        # Add spacing between bullet levels

    _add_page_number(slide, page_info[1], page_info[0], template)
    return slide


def build_two_column_slide(prs: Presentation, template: Template, title: str,
                           left_items: list[str], right_items: list[str],
                           page_info: tuple = (1, 1)):
    """Build a two-column comparison slide using Two Content layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_COL])
    _set_bg(slide, template.colors.bg)

    # Top accent line
    _add_shape(slide, MARGIN_LEFT, Inches(0.08), Inches(2.5), Inches(0.05),
               template.colors.accent)

    # Title
    _add_textbox(
        slide, MARGIN_LEFT, TITLE_TOP, CONTENT_WIDTH, Inches(0.8),
        title, template.fonts.title, 30, template.colors.accent, bold=True,
    )

    # Separator line under title
    _add_shape(slide, MARGIN_LEFT, Inches(1.2), CONTENT_WIDTH, Inches(0.02),
               template.colors.highlight)

    # Column headers and content
    col_width = Inches(5.0)
    col_top = Inches(1.6)
    col_height = Inches(5.0)

    # Left column
    _add_textbox(
        slide, MARGIN_LEFT, Inches(1.4), col_width, Inches(0.4),
        "●", template.fonts.body, 12, template.colors.accent, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    left_tf = _add_textbox(
        slide, MARGIN_LEFT, col_top, col_width, col_height,
        "", template.fonts.body, 16, template.colors.text,
    )
    if left_items:
        left_tf.paragraphs[0].text = left_items[0]
        if len(left_items) > 1:
            _add_bullets(left_tf, left_items[1:], template.fonts.body, 16, template.colors.text)

    # Vertical divider
    div_x = Inches(6.5)
    _add_shape(slide, div_x, Inches(1.4), Inches(0.02), Inches(5.5),
               template.colors.highlight)

    # Right column
    _add_textbox(
        slide, Inches(6.8), Inches(1.4), col_width, Inches(0.4),
        "●", template.fonts.body, 12, template.colors.accent, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    right_tf = _add_textbox(
        slide, Inches(6.8), col_top, col_width, col_height,
        "", template.fonts.body, 16, template.colors.text,
    )
    if right_items:
        right_tf.paragraphs[0].text = right_items[0]
        if len(right_items) > 1:
            _add_bullets(right_tf, right_items[1:], template.fonts.body, 16, template.colors.text)

    _add_page_number(slide, page_info[1], page_info[0], template)
    return slide


def build_ending_slide(prs: Presentation, template: Template, title: str = "Thank You",
                       subtitle: str = "", page_info: tuple = (1, 1)):
    """Build an ending slide with highlight background."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _set_bg(slide, template.colors.highlight)

    # Large decorative circle (subtle, alpha stripped)
    _add_shape(slide, Inches(9), Inches(-1), Inches(5), Inches(5),
               "#E8E8E8")

    # Title
    _add_textbox(
        slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
        title, template.fonts.title, 44, template.colors.bg,
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
            subtitle, template.fonts.body, 24, "#F0F0F0",
            alignment=PP_ALIGN.CENTER,
        )

    # Contact info placeholder
    _add_textbox(
        slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
        "Generated by PPTGenius",
        template.fonts.body, 12, "#D0D0D0",
        alignment=PP_ALIGN.CENTER,
    )

    _add_page_number(slide, page_info[1], page_info[0], template)
    return slide
