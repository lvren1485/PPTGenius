from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt

from ppt_generator.outline.models import Outline


def export_outline_to_pptx(outline: Outline, path: str) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for slide in outline.slides:
        s = prs.slides.add_slide(blank_layout)
        title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.85))
        tf = title_box.text_frame
        tf.text = slide.title
        for p in tf.paragraphs:
            p.font.size = Pt(26)
            p.font.bold = True

        para_text = (slide.body_paragraph or "").strip()
        bullets = list(slide.bullets)
        if not para_text and bullets:
            para_text = bullets[0][:420].strip()
            bullets = bullets[1:]

        if para_text:
            para_shape = s.shapes.add_textbox(Inches(0.5), Inches(1.12), Inches(9), Inches(1.45))
            ptf = para_shape.text_frame
            ptf.word_wrap = True
            ptf.text = para_text
            for p in ptf.paragraphs:
                p.font.size = Pt(14)
                p.space_after = Pt(4)
            bullet_top = Inches(2.62)
            bullet_h = Inches(4.12)
        else:
            bullet_top = Inches(1.12)
            bullet_h = Inches(5.62)

        body = s.shapes.add_textbox(Inches(0.5), bullet_top, Inches(9), bullet_h)
        btf = body.text_frame
        btf.word_wrap = True
        lines: list[str] = []
        for b in bullets:
            if b.strip():
                lines.append(f"• {b.strip()}")
        if not lines and slide.speaker_notes:
            lines.append(slide.speaker_notes[:1200])
        btf.text = "\n".join(lines) if lines else "（本页正文待补充）"
        for p in btf.paragraphs:
            p.font.size = Pt(17)
            p.space_after = Pt(3)

        if slide.speaker_notes.strip():
            notes_slide = s.notes_slide
            notes_slide.notes_text_frame.text = slide.speaker_notes

    prs.save(path)
