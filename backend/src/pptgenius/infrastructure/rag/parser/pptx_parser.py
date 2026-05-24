""".pptx parser — extracts slide text via python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .base import ParsedDocument, _normalise_text


def parse_pptx(file_path: str | Path) -> ParsedDocument:
    """Extract text from each slide in a .pptx file."""
    path = Path(file_path)
    prs = Presentation(str(path))

    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                rows: list[str] = []
                for row in table.rows:
                    cells = " | ".join(cell.text.strip() for cell in row.cells)
                    rows.append(f"| {cells} |")

                if rows:
                    # Skip leading entirely-empty rows
                    while rows and all(c.strip() == "" for c in rows[0].strip("|").split("|")):
                        rows.pop(0)
                    if rows:
                        header_sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
                        rows.insert(1, header_sep)
                        parts.extend(rows)
        slides.append("\n".join(parts))

    raw = "\n\n".join(slides)
    text = _normalise_text(raw)
    return ParsedDocument(
        file_path=str(path),
        file_type="pptx",
        text=text,
        metadata={
            "filename": path.name,
            "slides": len(prs.slides),
        },
    )
