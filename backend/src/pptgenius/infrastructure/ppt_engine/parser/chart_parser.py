"""Chart element renderer.

Supports category charts (column/bar/line/pie/doughnut/area/radar),
scatter charts (XyChartData), and bubble charts (BubbleChartData).
"""

import logging

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_LABEL_POSITION,
)

from .base import ChartElement, ChartData, XyChartData as XyModel, BubbleChartData as BubbleModel
from .styles import set_chart_area_fill, set_plot_area_fill

logger = logging.getLogger("ppt_engine.chart")

# ── Category chart type map (column/bar/line/pie/doughnut/area/radar) ──
CATEGORY_CHART_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked": XL_CHART_TYPE.LINE_STACKED,
    "line_stacked_100": XL_CHART_TYPE.LINE_STACKED_100,
    "pie": XL_CHART_TYPE.PIE,
    "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut_exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "area_stacked_100": XL_CHART_TYPE.AREA_STACKED_100,
    "radar": XL_CHART_TYPE.RADAR,
    "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
    "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
}

SCATTER_CHART_MAP = {
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
    "scatter_lines_no_markers": XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
}

BUBBLE_CHART = XL_CHART_TYPE.BUBBLE

LEGEND_POSITION_MAP = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "right": XL_LEGEND_POSITION.RIGHT,
    "left": XL_LEGEND_POSITION.LEFT,
    "top": XL_LEGEND_POSITION.TOP,
}

LABEL_POSITION_MAP = {
    "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
    "inside_end": XL_LABEL_POSITION.INSIDE_END,
    "center": XL_LABEL_POSITION.CENTER,
    "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
    "above": XL_LABEL_POSITION.ABOVE,
    "best_fit": XL_LABEL_POSITION.BEST_FIT,
}


def render_chart(slide, el: ChartElement) -> None:
    """Render a chart element onto a slide. Dispatches by chart type."""
    left = Inches(el.position.left)
    top = Inches(el.position.top)
    width = Inches(el.position.width)
    height = Inches(el.position.height)

    logger.debug("Rendering chart: type=%s pos=(%.1f,%.1f)", el.chart_type, el.position.left, el.position.top)

    if el.is_xy:
        chip = _render_scatter(slide, el, left, top, width, height)
    elif el.is_bubble:
        chip = _render_bubble(slide, el, left, top, width, height)
    else:
        chip = _render_category(slide, el, left, top, width, height)

    _apply_chart_style(chip, el)


def _render_category(slide, el, left, top, width, height):
    """Render a category chart (column/bar/line/pie/area/radar)."""
    chart_type = CATEGORY_CHART_MAP[el.chart_type]
    chart_data = CategoryChartData()
    d: ChartData = el.data
    chart_data.categories = d.categories
    for s in d.series:
        chart_data.add_series(s.name, s.values)

    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    if el.title:
        chart.has_title = True
        tf = chart.chart_title.text_frame
        if tf.paragraphs:
            tf.paragraphs[0].text = el.title
    return chart


def _render_scatter(slide, el, left, top, width, height):
    """Render a scatter (XY) chart."""
    chart_type = SCATTER_CHART_MAP[el.chart_type]
    chart_data = XyChartData()
    d: XyModel = el.data
    for s in d.series:
        series = chart_data.add_series(s.name)
        for pt in s.points:
            series.add_data_point(pt.x, pt.y)

    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    if el.title:
        chart.has_title = True
        tf = chart.chart_title.text_frame
        if tf.paragraphs:
            tf.paragraphs[0].text = el.title
    return chart


def _render_bubble(slide, el, left, top, width, height):
    """Render a bubble chart."""
    chart_data = BubbleChartData()
    d: BubbleModel = el.data
    for s in d.series:
        series = chart_data.add_series(s.name)
        for pt in s.points:
            series.add_data_point(pt.x, pt.y, pt.size)

    graphic_frame = slide.shapes.add_chart(BUBBLE_CHART, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    if el.title:
        chart.has_title = True
        tf = chart.chart_title.text_frame
        if tf.paragraphs:
            tf.paragraphs[0].text = el.title
    return chart


def _apply_chart_style(chart, el: ChartElement) -> None:
    """Apply common style properties to any chart type."""
    style = el.style
    if not style:
        return

    if style.has_legend:
        chart.has_legend = True
        chart.legend.position = LEGEND_POSITION_MAP.get(style.legend_position, XL_LEGEND_POSITION.BOTTOM)
        chart.legend.include_in_layout = False

    if style.has_data_labels and chart.plots:
        plot = chart.plots[0]
        plot.has_data_labels = True
        if style.data_label_position:
            plot.data_labels.position = LABEL_POSITION_MAP.get(
                style.data_label_position, XL_LABEL_POSITION.OUTSIDE_END
            )

    if style.series_colors:
        # Pie/doughnut: colors apply to points (sectors) within the single series
        if el.chart_type.startswith("pie") or el.chart_type.startswith("doughnut"):
            series = chart.series[0]
            for i, c in enumerate(style.series_colors):
                if c == "none" or i >= len(series.points):
                    continue
                series.points[i].format.fill.solid()
                series.points[i].format.fill.fore_color.rgb = RGBColor.from_string(c)
        else:
            for i, c in enumerate(style.series_colors):
                if c == "none" or i >= len(chart.series):
                    continue
                chart.series[i].format.fill.solid()
                chart.series[i].format.fill.fore_color.rgb = RGBColor.from_string(c)

    # Chart area / plot area background
    if style.chart_area_fill:
        if style.chart_area_fill == "none":
            set_chart_area_fill(chart, no_fill=True)
        else:
            set_chart_area_fill(chart, color_hex=style.chart_area_fill)
    if style.plot_area_fill:
        if style.plot_area_fill == "none":
            set_plot_area_fill(chart, no_fill=True)
        else:
            set_plot_area_fill(chart, color_hex=style.plot_area_fill)

    # Font sizes
    if style.title_font_size and chart.has_title:
        tf = chart.chart_title.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].font.size = Pt(style.title_font_size)
    if style.axis_font_size:
        try:
            chart.category_axis.tick_labels.font.size = Pt(style.axis_font_size)
        except ValueError:
            pass
        try:
            chart.value_axis.tick_labels.font.size = Pt(style.axis_font_size)
        except ValueError:
            pass
