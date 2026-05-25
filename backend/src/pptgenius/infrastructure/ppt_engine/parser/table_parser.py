"""Table element renderer."""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .base import TableElement
from .styles import set_cell_fill_and_border


def render_table(slide, el: TableElement) -> None:
    """Render a table element onto a slide."""
    left = Inches(el.position.left)
    top = Inches(el.position.top)
    width = Inches(el.position.width)
    height = Inches(el.position.height) if el.position.height else Inches(1.0)

    shape = slide.shapes.add_table(el.rows, el.cols, left, top, width, height)
    table = shape.table

    # Column widths
    if el.col_widths:
        for i, w in enumerate(el.col_widths):
            if i < len(table.columns):
                table.columns[i].width = Inches(w)

    # Cell content
    for cell_spec in el.cells:
        cell = table.cell(cell_spec.row, cell_spec.col)
        cell.text = cell_spec.text
        if cell_spec.font:
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _apply_cell_font(r, cell_spec.font)

    # Merges
    if el.merges:
        for m in el.merges:
            c1 = table.cell(m.from_row, m.from_col)
            c2 = table.cell(m.to_row, m.to_col)
            c1.merge(c2)

    # Header style
    if el.header:
        h = el.header
        for c in range(el.cols):
            cell = table.cell(h.row, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(h.fill)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor.from_string(h.font_color)
                    if h.font_bold:
                        r.font.bold = True
                    if h.font_size:
                        r.font.size = Pt(h.font_size)
            # Header border
            border_color = h.fill.lstrip("#")
            set_cell_fill_and_border(cell, fill_hex=h.fill, border_hex="ffffff", border_width_pt=0.5)

    # Table-wide borders + stripe
    if el.style:
        for r in range(el.rows):
            for c in range(el.cols):
                cell = table.cell(r, c)
                if el.header and r == el.header.row:
                    continue  # header borders already set
                fill = None
                if el.style.stripe_even and r % 2 == 0:
                    fill = el.style.stripe_even
                elif el.style.stripe_odd and r % 2 == 1:
                    fill = el.style.stripe_odd
                hex_color = el.style.border_color.lstrip("#")
                set_cell_fill_and_border(
                    cell,
                    fill_hex=fill,
                    border_hex=hex_color,
                    border_width_pt=el.style.border_width_pt,
                )


def _apply_cell_font(run, font_spec) -> None:
    f = run.font
    if font_spec.name: f.name = font_spec.name
    if font_spec.size: f.size = Pt(font_spec.size)
    if font_spec.bold is not None: f.bold = font_spec.bold
    if font_spec.italic is not None: f.italic = font_spec.italic
    if font_spec.color: f.color.rgb = RGBColor.from_string(font_spec.color)
