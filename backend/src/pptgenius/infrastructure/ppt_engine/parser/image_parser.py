"""Image element renderer. SVG auto-converted to PNG before embedding.

Primary: resvg-py (Rust, pre-built wheels for all platforms, no system deps).
Fallback: cairosvg (needs Cairo lib: apt install libcairo2 on Linux).
"""

import logging
import os

from pptx.util import Inches

from .base import ImageElement

logger = logging.getLogger("ppt_engine.image")


def render_picture(slide, el: ImageElement, workspace_path: str = ".") -> None:
    """Render a picture element onto a slide.

    Supports:
      - path:  direct file path (PNG / SVG, resolved against workspace_path)
      - name+color:  Tabler icon lookup → colored SVG → PNG
    """
    left = Inches(el.position.left)
    top = Inches(el.position.top)

    if el.name:
        from ..icon_search import get_colored_svg
        img_path = get_colored_svg(el.name, el.color, workspace_path)
        img_path = _svg_to_png(os.path.join(workspace_path, img_path))
    else:
        img_path = el.path
        if not os.path.isabs(img_path):
            img_path = os.path.join(workspace_path, img_path)

    if not os.path.exists(img_path):
        logger.warning("Image not found, inserting placeholder: %s", img_path)
        _render_image_placeholder(slide, el, img_path)
        return

    if img_path.lower().endswith(".svg"):
        img_path = _svg_to_png(img_path)

    if el.fit == "aspect":
        width = Inches(el.position.width)
        if el.position.height is not None:
            slide.shapes.add_picture(img_path, left, top, width, Inches(el.position.height))
        else:
            slide.shapes.add_picture(img_path, left, top, width)
    else:
        width = Inches(el.position.width)
        height = Inches(el.position.height) if el.position.height else width
        slide.shapes.add_picture(img_path, left, top, width, height)


def _render_image_placeholder(slide, el: ImageElement, img_path: str) -> None:
    """Insert a placeholder shape when the image file doesn't exist.

    Creates a rounded rectangle with a border, containing the image filename
    as text so the user can see what was supposed to be there.
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    left = Inches(el.position.left)
    top = Inches(el.position.top)
    width = Inches(el.position.width)
    height = Inches(el.position.height) if el.position.height else width

    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[图片]\n{os.path.basename(img_path)}"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x96, 0x96, 0x96)

    from pptx.oxml.ns import qn
    txBody = shape._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def _svg_to_png(svg_path: str) -> str:
    """Convert SVG to PNG at high resolution (min 500px per dimension).

    Renders at a target width derived from the SVG viewBox so the output
    is always >= 500px on the shorter side.  Resvg primary, cairosvg fallback.
    """
    png_path = svg_path.rsplit(".", 1)[0] + ".png"

    # Cache: skip re-conversion if PNG is newer than SVG source
    if os.path.exists(png_path) and os.path.getmtime(png_path) >= os.path.getmtime(svg_path):
        return png_path

    target_w, target_h = _calc_render_size(svg_path)

    # Primary: resvg-py (pre-built wheels, zero system deps)
    try:
        import resvg_py
        png_bytes = resvg_py.svg_to_bytes(
            svg_path=svg_path, width=target_w, height=target_h,
        )
        with open(png_path, "wb") as f:
            f.write(png_bytes)
        logger.info("SVG → PNG (resvg, %dx%d): %s", target_w, target_h, os.path.basename(svg_path))
        return png_path
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("resvg-py failed: %s, trying cairosvg fallback", exc)

    # Fallback: cairosvg (needs libcairo2 on Linux, GTK runtime on Windows)
    try:
        import cairosvg
        cairosvg.svg2png(
            url=svg_path, write_to=png_path,
            output_width=target_w, output_height=target_h,
        )
        logger.info("SVG → PNG (cairosvg, %dx%d): %s", target_w, target_h, os.path.basename(svg_path))
        return png_path
    except ImportError:
        logger.warning("No SVG converter available. Install resvg-py or cairosvg.")
        return svg_path
    except OSError as exc:
        logger.warning("Cairo library not found (%s). Install libcairo2 or use resvg-py.", exc)
        return svg_path


def _calc_render_size(svg_path: str) -> tuple[int, int]:
    """Calculate render size from SVG viewBox so the shorter side >= 500px."""
    import re

    target_min = 1024
    try:
        with open(svg_path, "r", encoding="utf-8") as f:
            head = f.read(4096)
    except Exception:
        return target_min, target_min

    m = re.search(r'viewBox\s*=\s*["\']([^"\']+)["\']', head)
    if not m:
        return target_min, target_min

    parts = m.group(1).split()
    if len(parts) < 4:
        return target_min, target_min

    try:
        vb_w, vb_h = float(parts[2]), float(parts[3])
    except ValueError:
        return target_min, target_min

    if vb_w <= 0 or vb_h <= 0:
        return target_min, target_min

    # Scale so the shorter side reaches target_min
    if vb_w >= vb_h:
        target_w = target_min
        target_h = max(1, round(target_min * vb_h / vb_w))
    else:
        target_h = target_min
        target_w = max(1, round(target_min * vb_w / vb_h))

    return target_w, target_h
