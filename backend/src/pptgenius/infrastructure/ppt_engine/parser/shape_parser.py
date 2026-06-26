"""Shape element renderer."""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .base import ShapeElement
from .shapes import resolve_shape
from .styles import (
    apply_shape_fill,
    apply_shape_line,
    set_shape_picture_fill,
    add_gradient_stop,
    set_run_fonts,
)


def render_shape(slide, el: ShapeElement) -> None:
    """Render an AutoShape element onto a slide."""
    left = Inches(el.position.left)
    top = Inches(el.position.top)
    width = Inches(el.position.width)
    height = Inches(el.position.height) if el.position.height else Inches(1.0)

    mso_shape = resolve_shape(el.shape_type)
    if mso_shape is None:
        raise ValueError(f"unknown shape_type '{el.shape_type}'")
    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

    # Fill
    if el.fill:
        if el.fill.type == "solid":
            apply_shape_fill(shape, "solid", el.fill.color)
        elif el.fill.type == "no_fill":
            apply_shape_fill(shape, "no_fill")
        elif el.fill.type == "gradient":
            apply_shape_fill(shape, "gradient")
            if el.fill.gradient_angle is not None:
                shape.fill.gradient_angle = el.fill.gradient_angle
            if el.fill.gradient_stops:
                for gs in el.fill.gradient_stops:
                    add_gradient_stop(shape, gs.position, gs.color.lstrip("#"))
        elif el.fill.type == "picture":
            set_shape_picture_fill(shape, el.fill.color or "", "stretch")

    # Line
    if el.line and el.line.color:
        apply_shape_line(shape, el.line.color, el.line.width_pt, el.line.dash)

    # Rotation
    if el.rotation:
        shape.rotation = el.rotation

    # Text inside shape
    if el.text:
        tf = shape.text_frame
        tf.word_wrap = True
        for i, block in enumerate(el.text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _apply_paragraph(p, block.paragraph)


ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _apply_paragraph(p, para_spec) -> None:
    p.alignment = ALIGN_MAP.get(para_spec.alignment, PP_ALIGN.LEFT)
    p.level = para_spec.level
    if para_spec.space_before_pt is not None:
        p.space_before = Pt(para_spec.space_before_pt)
    if para_spec.space_after_pt is not None:
        p.space_after = Pt(para_spec.space_after_pt)
    for run_spec in para_spec.runs:
        run = p.add_run()
        run.text = run_spec.text
        if run_spec.font:
            from pptx.dml.color import RGBColor
            f = run.font
            s = run_spec.font
            f.name = s.name or "微软雅黑"; set_run_fonts(run, f.name)
            if s.size: f.size = Pt(s.size)
            if s.bold is not None: f.bold = s.bold
            if s.color: f.color.rgb = RGBColor.from_string(s.color)
