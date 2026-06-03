"""PPT Generator — validates then builds .pptx from instruction JSON."""

from __future__ import annotations

import logging
import os
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from .parser.base import PPTInstruction, resolve_position
from .validator import ValidationResult, validate_instruction
from .parser.chart_parser import render_chart
from .parser.table_parser import render_table
from .parser.text_parser import render_textbox
from .parser.image_parser import render_picture
from .parser.shape_parser import render_shape

logger = logging.getLogger("ppt_engine.generator")

# Built-in parent bounds (slide, common columns)
SLIDE_BOUNDS = ("slide", 0, 0, 13.333, 7.5)


async def generate_ppt(
    data: dict[str, Any],
    output_path: str,
    workspace_path: str | None = None,  # A path only for temp files without any important data
) -> dict[str, Any]:
    """Validate instruction JSON, then generate .pptx.

    Returns:
        {"ok": True, "path": "...", "slide_count": N} or {"ok": False, "errors": [...]}
    """
    logger.info("Starting PPT generation → %s", output_path)

    # ── Validate ──
    result: ValidationResult = validate_instruction(data)
    if not result.is_valid:
        logger.warning("Validation failed, aborting generation")
        return {"ok": False, "errors": result.errors}

    instruction: PPTInstruction = result.parsed
    meta = instruction.meta
    logger.info("Validation OK, building %d slide(s)", len(instruction.slides))

    # ── Build ──
    prs = Presentation()
    prs.slide_width = Inches(meta.slide_width)
    prs.slide_height = Inches(meta.slide_height)

    # Fix sldSz type attribute — python-pptx default template is 4:3,
    # leaving type="screen4x3" even when dimensions are 16:9.  PowerPoint
    # shows a repair prompt when type and actual size disagree.
    _sldSz = prs.element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz'
    )
    if _sldSz is not None and _sldSz.get('type') != 'screen16x9':
        _sldSz.set('type', 'screen16x9')

    ws = workspace_path or "."

    for si, slide_spec in enumerate(instruction.slides):
        layout_idx = _find_layout_index(prs, slide_spec.layout)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        logger.debug("Slide %d/%d: layout=%s (%d element(s))",
                     si + 1, len(instruction.slides), slide_spec.layout, len(slide_spec.elements))

        if slide_spec.background:
            _apply_slide_background(slide, slide_spec.background)

        # Resolve relative positions: scan for container shapes, build parent_bounds
        parent_bounds = _build_parent_bounds(slide_spec, meta)

        for ei, element in enumerate(slide_spec.elements):
            # Apply relative position resolution
            if element.position.parent:
                element.position = resolve_position(element.position, parent_bounds)

            try:
                if element.type == "textbox":
                    render_textbox(slide, element)
                elif element.type == "chart":
                    render_chart(slide, element)
                elif element.type == "table":
                    render_table(slide, element)
                elif element.type == "picture":
                    render_picture(slide, element, ws)
                elif element.type == "shape":
                    render_shape(slide, element)
            except Exception as exc:
                logger.exception("Render failed: slides[%d].elements[%d] (%s)", si, ei, element.type)
                return {"ok": False, "errors": [{
                    "path": f"slides[{si}].elements[{ei}]",
                    "error": f"Render failed: {exc}",
                }]}

        # Slide notes
        if slide_spec.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_spec.notes
            logger.debug("Notes added to slide %d", si + 1)

    # ── Save ──
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    file_size = os.path.getsize(output_path)
    logger.info("PPT saved: %s (%d bytes, %d slides)", output_path, file_size, len(instruction.slides))

    # Clean up temp icon files (colored SVGs + PNGs generated per-run)
    from .icon_search import cleanup_temp_icons
    cleanup_temp_icons(workspace_path)

    return {"ok": True, "path": output_path, "slide_count": len(instruction.slides), "file_size": file_size}


def _build_parent_bounds(slide_spec, meta) -> dict[str, tuple[float, float, float, float]]:
    """Build parent bounds from container shapes on the slide.

    A shape without a parent that has an explicit fill acts as a container.
    The container name is inferred from its position pattern.
    Returns dict of parent_id → (left, top, width, height).
    """
    bounds = {"slide": SLIDE_BOUNDS[1:]}
    for el in slide_spec.elements:
        if el.type == "shape" and not el.position.parent and el.fill and el.fill.type in ("solid", "gradient"):
            # Determine container name from position pattern
            l, t, w, h = el.position.left, el.position.top, el.position.width, el.position.height or 0
            if l < meta.slide_width / 2:
                name = "left_col"
            else:
                name = "right_col"
            bounds[name] = (l, t, w, h)
    return bounds


def _apply_slide_background(slide, bg) -> None:
    """Apply slide background (solid color / gradient / image)."""
    from pptx.dml.color import RGBColor
    fill = slide.background.fill

    if bg.type == "solid" and bg.color:
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg.color)
    elif bg.type == "gradient":
        fill.gradient()
        if bg.gradient_angle is not None:
            fill.gradient_angle = bg.gradient_angle
        if bg.gradient_stops:
            from .parser.styles import add_gradient_stop
            # Replace default stops with custom ones
            for i, gs in enumerate(bg.gradient_stops):
                if i < 2:
                    # Modify existing stops
                    stops = fill.gradient_stops
                    if i == 0:
                        stops[0].color.rgb = RGBColor.from_string(gs.color)
                    elif i == 1:
                        stops[1].color.rgb = RGBColor.from_string(gs.color)
                else:
                    # Add extra stops via lxml
                    _shape = type("obj", (), {"fill": fill})()
                    add_gradient_stop(_shape, gs.position, gs.color.lstrip("#"))
    elif bg.type == "image" and bg.color:
        from .parser.styles import set_slide_image_background
        import os as _os
        img_path = bg.color
        if not _os.path.isabs(img_path):
            img_path = _os.path.abspath(_os.path.join(_os.getcwd(), img_path))
        set_slide_image_background(slide, img_path)


def _find_layout_index(prs, layout_name: str) -> int:
    """Find layout index by name. Falls back to blank (6)."""
    for i, lo in enumerate(prs.slide_layouts):
        if lo.name == layout_name:
            return i
    for i, lo in enumerate(prs.slide_layouts):
        if layout_name in lo.name or lo.name in layout_name:
            return i
    return 6
