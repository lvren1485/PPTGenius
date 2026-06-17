"""Generate a sample PPT to verify Chinese font rendering (ea/cs typeface fix).

Run:  python -m src.tests.test_font_export
Output:  data/test_output/font_test.pptx
"""

from __future__ import annotations

import asyncio
import json
import sys
from pathlib import Path

# Ensure backend/src is on sys.path
_REPO = Path(__file__).resolve().parent.parent.parent.parent
_SRC = _REPO / "backend" / "src"
if str(_SRC) not in sys.path:
    sys.path.insert(0, str(_SRC))

from pptgenius.infrastructure.ppt_engine.generator import generate_ppt

OUTPUT_DIR = _REPO / "backend" / "data" / "test_output"


def build_instruction() -> dict:
    """Build a 5-slide PPT instruction with Chinese fonts in every text element."""
    return {
        "meta": {
            "slide_width": 13.333,
            "slide_height": 7.5,
            "language": "zh",
        },
        "slides": [
            # ── Slide 0: Title ──
            {
                "layout": "title_slide",
                "background": {"type": "solid", "color": "1a73e8"},
                "elements": [
                    {
                        "type": "textbox",
                        "position": {"left": 1.5, "top": 2.0, "width": 10.3, "height": 2.0},
                        "content": [{
                            "paragraph": {
                                "alignment": "center",
                                "runs": [{
                                    "text": "PPTGenius 字体测试报告",
                                    "font": {"name": "微软雅黑", "size": 40, "bold": True, "color": "ffffff"},
                                }],
                            }
                        }],
                    },
                    {
                        "type": "textbox",
                        "position": {"left": 1.5, "top": 4.2, "width": 10.3, "height": 1.0},
                        "content": [{
                            "paragraph": {
                                "alignment": "center",
                                "runs": [{
                                    "text": "2024年6月 · AI驱动PPT生成平台",
                                    "font": {"name": "微软雅黑", "size": 18, "bold": False, "color": "c8d6e5"},
                                }],
                            }
                        }],
                    },
                ],
            },
            # ── Slide 1: 中文字体渲染对比 ──
            {
                "layout": "content_bullet",
                "background": {"type": "solid", "color": "f8f9fa"},
                "elements": [
                    {
                        "type": "textbox",
                        "position": {"left": 0.8, "top": 0.4, "width": 11.7, "height": 0.8},
                        "content": [{
                            "paragraph": {
                                "alignment": "left",
                                "runs": [{
                                    "text": "中文字体在多层级文本中的渲染验证",
                                    "font": {"name": "微软雅黑", "size": 30, "bold": True, "color": "1a73e8"},
                                }],
                            }
                        }],
                    },
                    # Large body text
                    {
                        "type": "textbox",
                        "position": {"left": 0.8, "top": 1.5, "width": 5.5, "height": 5.0},
                        "content": [
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 12,
                                    "runs": [{
                                        "text": "H1 一级标题 · 微软雅黑 36pt Bold",
                                        "font": {"name": "微软雅黑", "size": 36, "bold": True, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 8,
                                    "runs": [{
                                        "text": "H2 二级标题 · 28pt Bold — 展示粗体效果与较大字号",
                                        "font": {"name": "微软雅黑", "size": 28, "bold": True, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 8,
                                    "runs": [{
                                        "text": "H3 三级标题 · 22pt Bold — 进一步细化主题层次",
                                        "font": {"name": "微软雅黑", "size": 22, "bold": True, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 8,
                                    "runs": [{
                                        "text": "正文 Body · 16pt Regular — 这是标准的正文段落。",
                                        "font": {"name": "微软雅黑", "size": 16, "bold": False, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 0,
                                    "runs": [{
                                        "text": "说明 Caption · 14pt Regular — 脚注或数据来源说明文字。",
                                        "font": {"name": "微软雅黑", "size": 14, "bold": False, "color": "5f6368"},
                                    }],
                                }
                            },
                        ],
                    },
                    # Right column: mixed formatting demo
                    {
                        "type": "textbox",
                        "position": {"left": 7.0, "top": 1.5, "width": 5.5, "height": 5.0},
                        "content": [
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 12,
                                    "runs": [{
                                        "text": "中英混排 Mixed CJK+Latin",
                                        "font": {"name": "微软雅黑", "size": 24, "bold": True, "color": "1a73e8"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 10,
                                    "runs": [
                                        {"text": "中文正文 + ", "font": {"name": "微软雅黑", "size": 16, "color": "202124"}},
                                        {"text": "English Text", "font": {"name": "微软雅黑", "size": 16, "bold": True, "italic": True, "color": "ea4335"}},
                                        {"text": " 混排在同一段落中，", "font": {"name": "微软雅黑", "size": 16, "color": "202124"}},
                                    ],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 10,
                                    "runs": [
                                        {"text": "测试", "font": {"name": "微软雅黑", "size": 22, "bold": True, "color": "202124"}},
                                        {"text": "不同字号与", "font": {"name": "微软雅黑", "size": 16, "color": "202124"}},
                                        {"text": "颜色", "font": {"name": "微软雅黑", "size": 22, "bold": True, "color": "34a853"}},
                                        {"text": "混排效果。", "font": {"name": "微软雅黑", "size": 16, "color": "202124"}},
                                    ],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 10,
                                    "runs": [{
                                        "text": "斜体 + 下划线 + 粗体 组合效果演示",
                                        "font": {"name": "微软雅黑", "size": 18, "bold": True, "italic": True, "underline": "single", "color": "ea4335"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 0,
                                    "runs": [{
                                        "text": "THE QUICK BROWN FOX 1234567890 !@#$%",
                                        "font": {"name": "微软雅黑", "size": 14, "bold": False, "color": "5f6368"},
                                    }],
                                }
                            },
                        ],
                    },
                ],
            },
            # ── Slide 2: 表格 + 中文 ──
            {
                "layout": "content_table",
                "background": {"type": "solid", "color": "ffffff"},
                "elements": [
                    {
                        "type": "textbox",
                        "position": {"left": 0.8, "top": 0.4, "width": 11.7, "height": 0.8},
                        "content": [{
                            "paragraph": {
                                "alignment": "left",
                                "runs": [{
                                    "text": "表格中的中文字体验证",
                                    "font": {"name": "微软雅黑", "size": 28, "bold": True, "color": "1a73e8"},
                                }],
                            }
                        }],
                    },
                    {
                        "type": "table",
                        "position": {"left": 0.8, "top": 1.6, "width": 11.7, "height": 4.5},
                        "rows": 5,
                        "cols": 4,
                        "col_widths": [2.5, 3.5, 2.8, 2.9],
                        "header": {"row": 0, "fill": "1a73e8", "font_color": "ffffff", "font_bold": True, "font_size": 16},
                        "style": {"border_color": "dadce0", "border_width_pt": 0.5, "stripe_even": "f0f4ff"},
                        "cells": [
                            {"row": 0, "col": 0, "text": "功能模块", "font": {"name": "微软雅黑", "size": 16, "bold": True, "color": "ffffff"}},
                            {"row": 0, "col": 1, "text": "技术方案", "font": {"name": "微软雅黑", "size": 16, "bold": True, "color": "ffffff"}},
                            {"row": 0, "col": 2, "text": "状态", "font": {"name": "微软雅黑", "size": 16, "bold": True, "color": "ffffff"}},
                            {"row": 0, "col": 3, "text": "完成日期", "font": {"name": "微软雅黑", "size": 16, "bold": True, "color": "ffffff"}},
                            {"row": 1, "col": 0, "text": "大纲生成", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 1, "col": 1, "text": "LLM + RAG 知识检索增强生成", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 1, "col": 2, "text": "已完成", "font": {"name": "微软雅黑", "size": 14, "color": "34a853"}},
                            {"row": 1, "col": 3, "text": "2024-05-20", "font": {"name": "微软雅黑", "size": 14, "color": "5f6368"}},
                            {"row": 2, "col": 0, "text": "风格选择", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 2, "col": 1, "text": "AI Agent 自适应配色 + 布局匹配", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 2, "col": 2, "text": "已完成", "font": {"name": "微软雅黑", "size": 14, "color": "34a853"}},
                            {"row": 2, "col": 3, "text": "2024-05-22", "font": {"name": "微软雅黑", "size": 14, "color": "5f6368"}},
                            {"row": 3, "col": 0, "text": "PPT 生成引擎", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 3, "col": 1, "text": "python-pptx + 结构化指令 JSON", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 3, "col": 2, "text": "迭代中", "font": {"name": "微软雅黑", "size": 14, "color": "fbbc04"}},
                            {"row": 3, "col": 3, "text": "2024-06-15", "font": {"name": "微软雅黑", "size": 14, "color": "5f6368"}},
                            {"row": 4, "col": 0, "text": "字体渲染修复", "font": {"name": "微软雅黑", "size": 14, "bold": True, "color": "ea4335"}},
                            {"row": 4, "col": 1, "text": "lxml 注入 a:ea / a:cs typeface 属性", "font": {"name": "微软雅黑", "size": 14, "color": "202124"}},
                            {"row": 4, "col": 2, "text": "本次修复", "font": {"name": "微软雅黑", "size": 14, "bold": True, "color": "ea4335"}},
                            {"row": 4, "col": 3, "text": "2024-06-16", "font": {"name": "微软雅黑", "size": 14, "color": "5f6368"}},
                        ],
                    },
                ],
            },
            # ── Slide 3: 图表 + 中文标题 ──
            {
                "layout": "content_chart",
                "background": {"type": "solid", "color": "f8f9fa"},
                "elements": [
                    {
                        "type": "textbox",
                        "position": {"left": 0.8, "top": 0.4, "width": 11.7, "height": 0.7},
                        "content": [{
                            "paragraph": {
                                "alignment": "left",
                                "runs": [{
                                    "text": "图表中的中文字体渲染",
                                    "font": {"name": "微软雅黑", "size": 28, "bold": True, "color": "1a73e8"},
                                }],
                            }
                        }],
                    },
                    {
                        "type": "chart",
                        "chart_type": "column_clustered",
                        "position": {"left": 0.8, "top": 1.5, "width": 7.5, "height": 5.0},
                        "title": "各季度营收对比（单位：百万元）",
                        "data": {
                            "categories": ["2024Q1", "2024Q2", "2024Q3", "2024Q4"],
                            "series": [
                                {"name": "产品A", "values": [120, 145, 168, 192]},
                                {"name": "产品B", "values": [85, 98, 112, 130]},
                                {"name": "产品C", "values": [65, 72, 68, 75]},
                            ],
                        },
                        "style": {
                            "has_legend": True,
                            "legend_position": "bottom",
                            "series_colors": ["1a73e8", "34a853", "fbbc04"],
                            "chart_area_fill": "ffffff",
                            "title_font_size": 16,
                            "axis_font_size": 12,
                        },
                    },
                    {
                        "type": "textbox",
                        "position": {"left": 8.8, "top": 1.5, "width": 3.7, "height": 5.0},
                        "content": [
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 10,
                                    "runs": [{
                                        "text": "关键发现",
                                        "font": {"name": "微软雅黑", "size": 20, "bold": True, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 8,
                                    "runs": [{
                                        "text": "1. 产品A全年增长60%，Q4达到峰值1.92亿。",
                                        "font": {"name": "微软雅黑", "size": 14, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 8,
                                    "runs": [{
                                        "text": "2. 产品B稳步增长，年度复合增长率约15%。",
                                        "font": {"name": "微软雅黑", "size": 14, "color": "202124"},
                                    }],
                                }
                            },
                            {
                                "paragraph": {
                                    "alignment": "left",
                                    "space_after_pt": 0,
                                    "runs": [{
                                        "text": "3. 产品C表现平稳，建议加大研发投入以提升竞争力。",
                                        "font": {"name": "微软雅黑", "size": 14, "color": "202124"},
                                    }],
                                }
                            },
                        ],
                    },
                ],
            },
            # ── Slide 4: 装饰形状 + 中文文本 ──
            {
                "layout": "ending",
                "background": {"type": "solid", "color": "1a73e8"},
                "elements": [
                    {
                        "type": "textbox",
                        "position": {"left": 1.5, "top": 2.5, "width": 10.3, "height": 1.5},
                        "content": [{
                            "paragraph": {
                                "alignment": "center",
                                "runs": [{
                                    "text": "感谢观看",
                                    "font": {"name": "微软雅黑", "size": 42, "bold": True, "color": "ffffff"},
                                }],
                            }
                        }],
                    },
                    {
                        "type": "textbox",
                        "position": {"left": 1.5, "top": 4.2, "width": 10.3, "height": 0.8},
                        "content": [{
                            "paragraph": {
                                "alignment": "center",
                                "runs": [{
                                    "text": "PPTGenius · AI 驱动 PPT 生成平台",
                                    "font": {"name": "微软雅黑", "size": 18, "bold": False, "color": "c8d6e5"},
                                }],
                            }
                        }],
                    },
                    {
                        "type": "textbox",
                        "position": {"left": 1.5, "top": 5.2, "width": 10.3, "height": 0.6},
                        "content": [{
                            "paragraph": {
                                "alignment": "center",
                                "runs": [{
                                    "text": "字体修复版本 · 2024-06-16 · 思源黑体 / 微软雅黑 / 等线",
                                    "font": {"name": "微软雅黑", "size": 14, "bold": False, "color": "96b8e0"},
                                }],
                            }
                        }],
                    },
                ],
            },
        ],
    }


async def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    instruction = build_instruction()
    output_path = str(OUTPUT_DIR / "font_test.pptx")

    print(f"Generating PPT with {len(instruction['slides'])} slides...")
    result = await generate_ppt(instruction, output_path)

    if result["ok"]:
        print(f"OK: {result['path']}")
        print(f"Slides: {result['slide_count']}, Size: {result['file_size']:,} bytes")
        # Dump XML snippet of first text run for manual inspection
        _inspect_first_run(output_path)
    else:
        print(f"FAILED: {json.dumps(result['errors'], ensure_ascii=False, indent=2)}")
        sys.exit(1)


def _inspect_first_run(pptx_path: str) -> None:
    """Print the <a:rPr> XML of the first text run for manual verification."""
    from pptx import Presentation
    from lxml import etree

    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                        )
                        if rPr is not None:
                            xml_str = etree.tostring(rPr, encoding="unicode", pretty_print=True)
                            print("\n── First run <a:rPr> (verify ea/cs present) ──")
                            print(xml_str[:800])
                        return


if __name__ == "__main__":
    asyncio.run(main())
